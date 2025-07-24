# -*- coding: utf-8 -*-

# --- [0. 비동기 문제 해결] ---
import nest_asyncio
nest_asyncio.apply()

import os
import sys
import re
import hashlib
from pathlib import Path
import json
from datetime import datetime, date

# --- [토큰 계산 라이브러리] ---
import tiktoken

import streamlit as st
from PyPDF2 import PdfReader
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_community.vectorstores import FAISS
from langchain_community.retrievers import BM25Retriever
from langchain.schema import Document
from langchain.chains import RetrievalQA
from langchain.prompts import PromptTemplate
from langchain.retrievers import EnsembleRetriever
import streamlit.components.v1 as components

# --------------------------------------------------------------------------
# [1. 일일 사용량 제한 설정]
# --------------------------------------------------------------------------
DAILY_TOKEN_LIMIT = 20_000
USAGE_FILE = Path(__file__).parent / "usage_data.json"

def load_usage_data():
    today_str = str(date.today())
    try:
        if USAGE_FILE.exists():
            with open(USAGE_FILE, 'r') as f: data = json.load(f)
            if data.get("date") == today_str: return data.get("tokens_used", 0)
    except (json.JSONDecodeError, FileNotFoundError): pass
    save_usage_data(0)
    return 0

def save_usage_data(tokens_used):
    with open(USAGE_FILE, 'w') as f:
        json.dump({"date": str(date.today()), "tokens_used": tokens_used}, f)

def count_tokens(text: str) -> int:
    try:
        encoding = tiktoken.get_encoding("cl100k_base")
        return len(encoding.encode(text, disallowed_special=()))
    except Exception: return len(text) // 2

# --------------------------------------------------------------------------
# [기본 유틸리티 함수 정의]
# --------------------------------------------------------------------------
if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")

def safe_unicode(text: str) -> str:
    if not isinstance(text, str): text = str(text)
    return text.encode("utf-8", errors="ignore").decode("utf-8", errors="ignore").strip()
def clean_text(text):
    if not isinstance(text, str): return ""
    text = text.replace("\x00", ""); text = re.sub(r"[-\u001F\u007F-\u009F]", "", text)
    text = re.sub(r"[\ud800-\udfff]", "", text)
    return text.encode("utf-8", errors="ignore").decode("utf-8", errors="ignore").strip()
def extract_text_from_pdf(path: Path) -> str:
    try:
        reader = PdfReader(str(path))
        return "\n".join(clean_text(page.extract_text() or "") for page in reader.pages)
    except Exception as e: st.warning(f"[PDF 추출 실패] {path.name}: {e}"); return ""
def extract_text_from_pptx(path: Path) -> str:
    try:
        prs = Presentation(str(path))
        return "\n".join(clean_text(shape.text) for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    except Exception as e: st.warning(f"[PPTX 추출 실패] {path.name}: {e}"); return ""
def extract_text_from_json(path: Path) -> list[Document]:
    documents = []
    try:
        with open(path, 'r', encoding='utf-8') as f: data = json.load(f)
        for item in data:
            if not isinstance(item, dict): continue
            item_type = item.get("type", "정보")
            item_name = item.get("name") or item.get("조합원_형태") or item.get("category")
            summary = f"유형: {item_type}, 이름: {item_name}\n"
            full_content = json.dumps(item, ensure_ascii=False, indent=2)
            documents.append(Document(page_content=summary + full_content, metadata={"source": str(path.name), "type": item_type, "name": item_name}))
        return documents
    except Exception as e: st.warning(f"[JSON 추출 실패] {path.name}: {e}"); return []
def compute_file_hash(file_paths):
    hash_md5 = hashlib.md5()
    for path in sorted(file_paths):
        with open(path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
    return hash_md5.hexdigest()

@st.cache_resource
def load_all_documents_with_hash(file_paths, file_hash):
    documents = []
    for path in file_paths:
        if path.suffix == ".json" or ("json" in path.name.lower() and path.suffix == ".txt"):
            json_docs = extract_text_from_json(path)
            for doc in json_docs:
                doc.page_content = "[가장 정확한 최신 정보 출처: JSON 데이터베이스]\n\n" + doc.page_content
            documents.extend(json_docs)
        elif path.suffix == ".pdf":
            text = extract_text_from_pdf(path)
            if text.strip():
                documents.append(Document(page_content="[참고 자료]\n\n" + text, metadata={"source": str(path.name)}))
        elif path.suffix == ".pptx":
            text = extract_text_from_pptx(path)
            if text.strip():
                documents.append(Document(page_content="[참고 자료]\n\n" + text, metadata={"source": str(path.name)}))
            
    if not documents:
        st.error("문서를 불러오지 못했습니다. data 폴더를 만들고 그 안에 문서들을 넣어주세요.")
        st.stop()
    return documents

@st.cache_resource
def split_documents_into_chunks(_documents):
    splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=300)
    return splitter.split_documents(_documents)

def preprocess_query(query):
    # 연관 검색어 사전
    synonyms = {
        "피플팀": ["피플팀", "피플팀 ER그룹"],
    }
    original_query = query.strip()
    expanded_terms = set([original_query])
    for keyword, related_terms in synonyms.items():
        if keyword in original_query:
            for term in related_terms:
                expanded_terms.add(term)
    if len(expanded_terms) > 1:
        additional_terms = expanded_terms - {original_query}
        expanded_query = f"{original_query} 또는 {' 또는 '.join(additional_terms)}"
    else:
        expanded_query = original_query
    if not expanded_query.endswith(("에 대해", "에 대한 정보", "?")):
         expanded_query += "에 대해"
    return expanded_query

# --- [수정된 QA_SYSTEM_PROMPT] ---
QA_SYSTEM_PROMPT = """
너는 제공된 문서들을 바탕으로 질문에 답하는 AI 상담사다. 답변은 반드시 다음 규칙에 따라 생성해야 한다.
1. 정보 우선순위 준수: 1순위 [JSON], 2순위 [PDF], 3순위 [PPTX] 순으로 정보를 신뢰하고 답변의 기준으로 삼아라.
2. 정보 충돌 시 해결: 문서 간 정보가 충돌하면, 반드시 더 높은 순위의 정보를 기준으로 답해야 한다. (JSON > PDF > PPTX)
3. 정보 부재 시 응답: 어떤 문서에서도 질문에 대한 명시적 언급을 찾을 수 없다면, "관련 정보를 찾을 수 없습니다."라고 명확히 답변해라. 추측해서 답변하지 마라.
4. 답변 형식: 사용자가 짧은 키워드로 질문하더라도, 관련된 정보를 찾아 최대한 상세하고 친절한 문장으로 완성해서 답변해라.
5. 정보가 부족할 경우: 질문에 대한 직접적인 설명이나 정의를 찾을 수 없지만, 관련 키워드가 언급된 내용이 있다면, '정의를 찾을 수 없지만, 다음과 같은 내용에서 언급되었습니다'라고 안내하며 해당 문맥을 인용해라.
"""
QA_QUESTION_PROMPT = PromptTemplate(template=QA_SYSTEM_PROMPT + "\n\n{context}\n\n입력: {question}\n답변:", input_variables=["context", "question"])

@st.cache_resource
def initialize_qa_chain(all_paths, api_key):
    file_hash = compute_file_hash(all_paths)
    docs = load_all_documents_with_hash(all_paths, file_hash)
    chunks = split_documents_into_chunks(docs)
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001", google_api_key=api_key)
    llm = ChatGoogleGenerativeAI(google_api_key=api_key, model="gemini-1.5-flash", temperature=0)
    bm25_retriever = BM25Retriever.from_documents(chunks)
    bm25_retriever.k = min(20, len(chunks))
    faiss_vectorstore = FAISS.from_documents(chunks, embeddings)
    faiss_retriever = faiss_vectorstore.as_retriever(search_kwargs={"k": 10})
    ensemble_retriever = EnsembleRetriever(retrievers=[bm25_retriever, faiss_retriever], weights=[0.6, 0.4])
    return RetrievalQA.from_chain_type(llm=llm, chain_type="stuff", retriever=ensemble_retriever, chain_type_kwargs={"prompt": QA_QUESTION_PROMPT}, return_source_documents=True)

# --- [앱 실행 부분] ---
google_api_key = os.getenv("GOOGLE_API_KEY", "").strip() or st.secrets.get("GOOGLE_API_KEY")
if not google_api_key:
    st.error("Google API 키가 설정되지 않았습니다. Hugging Face의 'Settings > Repository secrets'에 GOOGLE_API_KEY를 등록해주세요."); st.stop()

st.set_page_config(page_title="삼성전기 존중노조 AI 집사", layout="centered", page_icon="https://raw.githubusercontent.com/SHIN-HOON-SIK/semunion-chatbot/main/logo_union_hands.png")

# --- [최종 디자인 CSS 스타일] ---
st.markdown("""
<style>
    /* 질문 입력창 클릭 시 굵고 진한 파란색 테두리 적용 */
    div[data-testid="stChatInput"]:focus-within {
        border: 2px solid #005A9C !important;
        box-shadow: none !important;
    }
    /* 내부 입력창의 기본 포커스 효과는 완전히 제거 */
    div[data-testid="stChatInput"] textarea:focus,
    div[data-testid="stChatInput"] textarea:focus-visible {
        border: none !important;
        box-shadow: none !important;
        outline: none !important;
    }
    
    /* 사용자 아바타(아이콘) 숨기기 */
    div[data-testid="stChatMessage-user"] div[data-testid="stChatAvatar"] {
        display: none;
    }

    /* 사용자 메시지 컨테이너 오른쪽 정렬 */
    div[data-testid="stChatMessage-user"] {
        justify-content: flex-end;
    }

    /* 말풍선 공통 스타일 */
    div[data-testid="stMarkdownContainer"] {
        border-radius: 0.75rem;
        padding: 1rem;
        border: none;
        word-wrap: break-word;
    }
    
    /* 답변 말풍선 배경색 */
    div[data-testid="stChatMessage-assistant"] div[data-testid="stMarkdownContainer"] {
        background-color: #f1f3f5;
        color: #000;
    }

    /* 질문 말풍선 배경색 */
    div[data-testid="stChatMessage-user"] div[data-testid="stMarkdownContainer"] {
        background-color: #ffffff;
        border: 1px solid #e0e0e0;
        color: #000;
    }
</style>
""", unsafe_allow_html=True)

st.markdown("""
<div style='display: flex; align-items: center; justify-content: center; gap: 10px; margin-bottom: 20px;'>
    <img src="https://raw.githubusercontent.com/SHIN-HOON-SIK/semunion-chatbot/main/logo_union_hands.png" width="50"/>
    <h1 style='color: #0d1a44; margin: 0;'>삼성전기 존중노조 AI 집사</h1>
</div>
<p style='text-align: center;'>존중노조 집행부에서 등록한 자료를 기반으로 질문하신 내용에 답변해 드립니다.</p>
""", unsafe_allow_html=True)

# --- [메인 로직] ---
base_dir = Path(__file__).parent
data_dir = base_dir / "data"
if not data_dir.exists():
    st.error(f"'{data_dir.name}' 폴더를 찾을 수 없습니다."); st.stop()
doc_paths = list(data_dir.glob("**/*.*"))

qa_chain = None
if not doc_paths:
    st.warning("data 폴더에 문서 파일이 없습니다.")
else:
    try:
        qa_chain = initialize_qa_chain(all_paths=doc_paths, api_key=google_api_key)
    except Exception as e:
        st.error(f"앱 초기화 중 오류가 발생했습니다: {e}"); st.stop()

if "messages" not in st.session_state:
    st.session_state.messages = []
for message in st.session_state.messages:
    with st.chat_message(message["role"], avatar=message.get("avatar")):
        st.markdown(message["content"])

user_query = st.chat_input("궁금하신 내용은 여기에 질문을 해보세요!", key="user_input")
tokens_used_today = load_usage_data()

if user_query:
    # --- 질문 길이 제한 로직 추가 ---
    if len(user_query) > 20:
        st.warning("사용 요금 종량제 문제로 질문의 길이를 제한하고 있습니다. 죄송합니다.")
    else:
        st.session_state.messages.append({"role": "user", "content": user_query})
        with st.chat_message("user"):
            st.markdown(user_query)
        
        if tokens_used_today >= DAILY_TOKEN_LIMIT:
            limit_message = "죄송합니다. 해당 세션에서 답변해 드릴 수 있는 토큰을 모두 소진했습니다. 잠시 후 다시 시도해 주세요."
            st.warning(limit_message)
            st.session_state.messages.append({"role": "assistant", "content": limit_message, "avatar": "https://raw.githubusercontent.com/SHIN-HOON-SIK/semunion-chatbot/main/logo_union_hands.png"})
        elif qa_chain:
            with st.spinner("AI 집사가 등록된 자료를 훑어보고 있습니다... 👀⏳..."):
                try:
                    processed_query = preprocess_query(user_query)
                    result = qa_chain.invoke({"query": processed_query})
                    answer = result["result"]
                    
                    tokens_this_round = count_tokens(processed_query) + count_tokens(answer)
                    new_total_tokens = tokens_used_today + tokens_this_round
                    save_usage_data(new_total_tokens)
                    
                    assistant_avatar = "https://raw.githubusercontent.com/SHIN-HOON-SIK/semunion-chatbot/main/logo_union_hands.png"
                    response = "죄송하지만 관련된 내용을 찾을 수 없습니다. 공짜로 만든 챗봇이라 능력이 부족합니다. T.T 다르게 질문 부탁드립니다." if not answer or "문서에 해당 정보가 없습니다" in answer else answer
                    
                    st.session_state.messages.append({"role": "assistant", "content": response, "avatar": assistant_avatar})
                    with st.chat_message("assistant", avatar=assistant_avatar):
                        if "죄송하지만" in response: st.info(response)
                        else: st.success(response)
                    
                    st.markdown(f"<p style='text-align: right; font-size: 0.8em; color: grey;'>이번 응답: 약 {tokens_this_round:,} 토큰 | 세션 사용량: {new_total_tokens:,} / {DAILY_TOKEN_LIMIT:,} 토큰</p>", unsafe_allow_html=True)
                    
                    # --- [★ 스크롤 제어 스크립트 추가 ★] ---
                    scroll_script = """
                    <script>
                        window.setTimeout(function() {
                            const messages = window.parent.document.querySelectorAll('[data-testid="stChatMessage"]');
                            if (messages.length > 0) {
                                const lastMessage = messages[messages.length - 1];
                                lastMessage.scrollIntoView({ behavior: 'smooth', block: 'center', inline: 'nearest' });
                            }
                        }, 300);
                    </script>
                    """
                    components.html(scroll_script, height=0)

                except Exception as e:
                    st.error(f"답변 생성 중 오류가 발생했습니다: {e}")
        else:
            st.warning("data 폴더에 문서가 없어 답변을 생성할 수 없습니다.")

# --- [앱 푸터] ---
st.markdown("""
<hr style="margin-top: 3em; margin-bottom: 1em;">
<div style="text-align: center; font-size: 0.85em; color: gray;">
    수원시 영통구 매영로 159번길 19 광교 더 퍼스트 지식산업센터<br>
    사업자 등록번호 133-82-71927<br>
    위원장: 신훈식 | 대표번호: 010-9496-6517 | 이메일: hoonsik79@hanmail.net
</div>
""", unsafe_allow_html=True)
